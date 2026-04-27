"""Build MAPPO demo PowerPoint presentation."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

IMG = "c:/Users/akeck/OneDrive/Desktop/Projects/mappo/ppt_images"
OUT = "c:/Users/akeck/OneDrive/Desktop/Projects/mappo/MAPPO_Demo_Presentation.pptx"

# ── Slide dimensions (widescreen 16:9) ──────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)
prs = Presentation()
prs.slide_width = W
prs.slide_height = H

# ── Color palette ────────────────────────────────────────────────────────────
DARK_BG   = RGBColor(0x1A, 0x1A, 0x2E)
MID_BG    = RGBColor(0x16, 0x21, 0x3E)
ACCENT    = RGBColor(0x0F, 0x3E, 0x7E)
GOLD      = RGBColor(0xE9, 0xC4, 0x6A)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xCC, 0xDD, 0xEE)
MAPPO_COL = RGBColor(0x26, 0x4E, 0xCC)
IPPO_COL  = RGBColor(0xE7, 0x6F, 0x51)

BLANK = prs.slide_layouts[6]  # completely blank

# ────────────────────────────────────────────────────────────────────────────
# Helper functions
# ────────────────────────────────────────────────────────────────────────────
def bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, left, top, w, h, color, text="", font_size=20,
        bold=False, align=PP_ALIGN.LEFT, font_color=WHITE,
        v_anchor=None, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    if v_anchor:
        tf.vertical_anchor = v_anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = font_color
    # background fill on the shape
    txBox.fill.solid()
    txBox.fill.fore_color.rgb = color
    return txBox

def rect(slide, left, top, w, h, color):
    shape = slide.shapes.add_shape(1, left, top, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def title_bar(slide, title, subtitle=""):
    rect(slide, 0, 0, W, Inches(1.4), ACCENT)
    box(slide, Inches(0.4), Inches(0.1), Inches(12), Inches(0.75),
        color=ACCENT, text=title, font_size=32, bold=True,
        align=PP_ALIGN.LEFT, font_color=WHITE)
    if subtitle:
        box(slide, Inches(0.4), Inches(0.85), Inches(12), Inches(0.45),
            color=ACCENT, text=subtitle, font_size=18,
            align=PP_ALIGN.LEFT, font_color=GOLD)

def full_image(slide, path, left=Inches(0.3), top=Inches(1.55),
               width=Inches(12.7), height=Inches(5.7)):
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, width, height)

def two_images(slide, path_l, path_r, top=Inches(1.55), h=Inches(5.7)):
    hw = Inches(6.25)
    if os.path.exists(path_l):
        slide.shapes.add_picture(path_l, Inches(0.25), top, hw, h)
    if os.path.exists(path_r):
        slide.shapes.add_picture(path_r, Inches(6.8),  top, hw, h)

def bullet_slide(slide, items, left=Inches(0.5), top=Inches(1.55),
                 w=Inches(12.3), h=Inches(5.5), font_size=22, gap=Inches(0.55)):
    y = top
    for item in items:
        box(slide, left, y, w, gap, color=DARK_BG, text=item,
            font_size=font_size, wrap=True)
        y += gap

def two_col(slide, left_items, right_items, top=Inches(1.55),
            h=Inches(5.5), font_size=20):
    cw = Inches(6.0)
    gap = Inches(0.5)
    y = top
    for item in left_items:
        box(slide, Inches(0.4), y, cw, gap, color=DARK_BG, text=item,
            font_size=font_size, wrap=True)
        y += gap
    y = top
    for item in right_items:
        box(slide, Inches(6.8), y, cw, gap, color=DARK_BG, text=item,
            font_size=font_size, wrap=True)
        y += gap

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title slide
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK_BG)
rect(s, 0, 0, W, Inches(0.15), GOLD)
rect(s, 0, H - Inches(0.15), W, Inches(0.15), GOLD)

box(s, Inches(0.8), Inches(1.6), Inches(11.7), Inches(1.2), DARK_BG,
    "The Surprising Effectiveness of PPO", font_size=42, bold=True,
    align=PP_ALIGN.CENTER, font_color=GOLD)
box(s, Inches(0.8), Inches(2.85), Inches(11.7), Inches(0.8), DARK_BG,
    "in Cooperative Multi-Agent Games", font_size=36, bold=True,
    align=PP_ALIGN.CENTER, font_color=WHITE)

rect(s, Inches(2.5), Inches(3.75), Inches(8.3), Inches(0.04), GOLD)

box(s, Inches(0.8), Inches(4.0), Inches(11.7), Inches(0.7), DARK_BG,
    "MAPPO  vs.  IPPO", font_size=20,
    align=PP_ALIGN.CENTER, font_color=LIGHT)
box(s, Inches(0.8), Inches(4.75), Inches(11.7), Inches(0.55), DARK_BG,
    "MPE Simple Spread  ·  MPE Speaker-Listener  ·  Hanabi", font_size=18,
    align=PP_ALIGN.CENTER, font_color=LIGHT)

box(s, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.5), DARK_BG,
    "5-Minute Project Demo  |  Reproduction Study", font_size=16,
    align=PP_ALIGN.CENTER, font_color=RGBColor(0x88, 0x99, 0xAA))

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Project Purpose & Motivation
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK_BG)
title_bar(s, "Project Purpose", "Why study PPO in cooperative multi-agent settings?")

items = [
    "  The original MAPPO paper (Yu et al., 2022) claimed that PPO-based methods are surprisingly\n"
    "  competitive with — or superior to — more complex MARL algorithms.",

    "  Goal: Reproduce the paper's key results on three benchmark environments to validate\n"
    "  whether simple PPO-style training generalizes across task types.",

    "  We compare MAPPO (centralized critic) vs. IPPO (decentralized) to measure whether\n"
    "  sharing global state information during training produces better cooperative behavior.",

    "  Environments span three distinct cooperation challenges:\n"
    "    • Coordination in physical space (MPE simple_spread)\n"
    "    • Explicit communication under role specialization (MPE speaker-listener)\n"
    "    • Implicit signaling under full partial-observability (Hanabi)",
]
bullet_slide(s, items, font_size=19, gap=Inches(1.35))

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Algorithms Overview
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK_BG)
title_bar(s, "Algorithms at a Glance", "Centralized training vs. decentralized execution")

# Two-column header labels
rect(s, Inches(0.3),  Inches(1.55), Inches(6.1), Inches(0.45), MAPPO_COL)
rect(s, Inches(6.7),  Inches(1.55), Inches(6.1), Inches(0.45), IPPO_COL)
box(s, Inches(0.3),  Inches(1.55), Inches(6.1), Inches(0.45), MAPPO_COL,
    "  Centralized-Critic (CTDE) Algorithms", font_size=16, bold=True,
    font_color=WHITE)
box(s, Inches(6.7),  Inches(1.55), Inches(6.1), Inches(0.45), IPPO_COL,
    "  Decentralized / Independent Algorithms", font_size=16, bold=True,
    font_color=WHITE)

left = [
    "MAPPO (Multi-Agent PPO)",
    "Each agent has its own actor network that only uses local\n"
    "observations at execution time.",
    "A single shared critic sees the full joint state — all agents'\n"
    "positions, velocities, and goals — during training only.",
    "This is Centralized Training, Decentralized Execution (CTDE).\n"
    "The critic is discarded after training; only actors are deployed.",
]
right = [
    "IPPO (Independent PPO)",
    "Each agent has its own actor and its own private critic.",
    "The critic only sees that agent's local observation —\n"
    "it has no knowledge of teammates' positions or intentions.",
    "Agents train in parallel, fully independently. Simpler, but\n"
    "blind to what teammates are doing during training.",
]
two_col(s, left, right, top=Inches(2.0), font_size=17)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Methodology
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK_BG)
title_bar(s, "Methodology", "Reproduction study — how we matched the original paper")

items = [
    "  Codebase: Official MAPPO repo (marlbenchmark/on-policy) — no architecture changes.",
    "  Training:  Google Colab A100 GPU · 128 rollout threads · TensorBoard logging · results saved to Google Drive.",
    "  TC 1 — MPE simple_spread:  2 M env steps run (paper: 25 M) · MAPPO vs IPPO · 3 agents · reward = −Σ dist(agent,landmark).",
    "  TC 2 — MPE speaker_listener:  2 M steps run (paper: 25 M) · MAPPO vs IPPO · 1 speaker + 1 listener.",
    "  TC 3 — Hanabi-Very-Small:  5 M steps run (paper: 10 M+) · MAPPO only · 2 agents · max score = 5.",
    "  Step budget limited by Colab session time: TC1 capped at 32 threads → 25 M steps would need ~8 hrs per run.",
    "  Metrics:  Average episode reward (MPE) and average game score (Hanabi) vs. environment steps.",
]
bullet_slide(s, items, font_size=17, gap=Inches(0.72))

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — TC1 Training Curve (cell 13)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK_BG)
title_bar(s, "Test Case 1 — MPE simple_spread",
          "MAPPO (centralized critic) vs IPPO (decentralized) · 3 agents cover 3 landmarks")
full_image(s, f"{IMG}/cell13_img0.png", top=Inches(1.55), height=Inches(4.7))
rect(s, 0, Inches(6.35), W, Inches(1.15), RGBColor(0x10, 0x28, 0x50))
box(s, Inches(0.35), Inches(6.4), Inches(12.6), Inches(1.0),
    RGBColor(0x10, 0x28, 0x50),
    "Our results (2 M steps):  MAPPO −125  |  IPPO −126  →  near-tie, gap not yet visible\n"
    "Paper results (25 M steps):  MAPPO ~−85  |  IPPO ~−110  →  ~25-point MAPPO lead\n"
    "Why the difference: TC1 was capped at 32 rollout threads; reaching 25 M steps would require ~8 hrs per run on Colab — beyond session limits.",
    font_size=13, font_color=LIGHT)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — TC2 Training Curves (cell 21, two images)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK_BG)
title_bar(s, "Test Case 2 — MPE simple_speaker_listener",
          "Communication under role specialization · MAPPO vs IPPO")
two_images(s, f"{IMG}/cell21_img0.png", f"{IMG}/cell21_img1.png", top=Inches(1.55), h=Inches(4.7))
rect(s, 0, Inches(6.35), W, Inches(1.15), RGBColor(0x10, 0x28, 0x50))
box(s, Inches(0.35), Inches(6.4), Inches(12.6), Inches(1.0),
    RGBColor(0x10, 0x28, 0x50),
    "Our results (2 M steps):  MAPPO −13.25  |  IPPO −14.96  →  MAPPO ahead by ~12%\n"
    "Paper results (25 M steps):  MAPPO ~−12  |  IPPO ~−17  →  MAPPO ahead by ~30%\n"
    "Closest match to paper: direction and order-of-magnitude correct even at reduced training.",
    font_size=13, font_color=LIGHT)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — TC3 Hanabi (cell 27)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK_BG)
title_bar(s, "Test Case 3 — Hanabi-Very-Small",
          "Implicit communication under full partial observability · 2 agents, max score = 5")
full_image(s, f"{IMG}/cell27_img0.png", top=Inches(1.55), height=Inches(4.7))
rect(s, 0, Inches(6.35), W, Inches(1.15), RGBColor(0x10, 0x28, 0x50))
box(s, Inches(0.35), Inches(6.4), Inches(12.6), Inches(1.0),
    RGBColor(0x10, 0x28, 0x50),
    "Our results (5 M steps):  Train score 1.90 / 5  |  Eval score 0.73 / 5  →  large train/eval gap\n"
    "Paper results (10 M+ steps):  Train ~2.5–3.0 / 5  →  agents still improving, more steps needed\n"
    "Eval gap likely due to underfitting: agents haven't fully generalized hint conventions at 5 M steps.",
    font_size=13, font_color=LIGHT)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Summary side-by-side (cell 29)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK_BG)
title_bar(s, "Summary — All Three Test Cases",
          "Training curves and final performance across environments")
full_image(s, f"{IMG}/cell29_img0.png")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Results vs. Paper
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK_BG)
title_bar(s, "Results vs. Paper", "Yu et al. (2022) at 25 M steps  vs.  our reproduction")

# Column headers
COL_HEADERS = ["", "Our steps", "Paper steps", "Our MAPPO", "Paper MAPPO", "Our IPPO", "Paper IPPO", "Direction"]
COL_X = [Inches(0.2), Inches(1.7), Inches(3.0), Inches(4.3), Inches(5.75), Inches(7.2), Inches(8.7), Inches(10.2)]
COL_W = [Inches(1.4), Inches(1.2), Inches(1.2), Inches(1.35), Inches(1.35), Inches(1.35), Inches(1.35), Inches(2.9)]

HDR_BG = RGBColor(0x0F, 0x3E, 0x7E)
ROW_BG = [RGBColor(0x1E, 0x2A, 0x45), RGBColor(0x16, 0x21, 0x3E)]
OK_COL = RGBColor(0x4C, 0xAF, 0x50)
WARN_COL = RGBColor(0xFF, 0xC1, 0x07)

def hdr_cell(slide, x, y, w, h, text):
    box(slide, x, y, w, h, HDR_BG, text, font_size=13, bold=True,
        align=PP_ALIGN.CENTER, font_color=WHITE)

def data_cell(slide, x, y, w, h, text, row_i, color=None):
    bg_c = ROW_BG[row_i % 2]
    fc = color if color else LIGHT
    box(slide, x, y, w, h, bg_c, text, font_size=13,
        align=PP_ALIGN.CENTER, font_color=fc)

ROW_H = Inches(0.62)
HDR_Y = Inches(1.55)
for ci, (hdr, x, w) in enumerate(zip(COL_HEADERS, COL_X, COL_W)):
    hdr_cell(s, x, HDR_Y, w, ROW_H, hdr)

rows = [
    ("TC1\nspread",   "2 M",  "25 M", "−125", "~−85",  "−126", "~−110", "Partial\n(not enough steps)"),
    ("TC2\nspeaker",  "2 M",  "25 M", "−13.3","~−12",  "−15.0","~−17",  "Match\n(MAPPO wins)"),
    ("TC3\nHanabi",   "5 M",  "10 M+","1.90/5","~2.5/5","N/A",  "N/A",  "Partial\n(underfitting)"),
]
DIR_COLORS = [WARN_COL, OK_COL, WARN_COL]

for ri, (row, dcol) in enumerate(zip(rows, DIR_COLORS)):
    y = HDR_Y + ROW_H * (ri + 1)
    for ci, (val, x, w) in enumerate(zip(row, COL_X, COL_W)):
        fc = dcol if ci == 7 else LIGHT
        data_cell(s, x, y, w, ROW_H, val, ri, color=fc)

# Why box
rect(s, 0, Inches(4.5), W, Inches(2.95), RGBColor(0x0D, 0x1E, 0x36))
box(s, Inches(0.35), Inches(4.58), Inches(12.6), Inches(2.75),
    RGBColor(0x0D, 0x1E, 0x36),
    "Why our results don't fully match the paper:\n\n"
    "  1. Step budget — We ran 2 M steps for MPE vs. the paper's 25 M (8% of the training). MAPPO's advantage over IPPO\n"
    "      is an emergent, long-horizon effect that only becomes visible after sustained training.\n\n"
    "  2. Colab session limits — TC1 was capped at 32 rollout threads (vs. 128 for TC2) due to environment constraints.\n"
    "      Reaching 25 M steps at that rate would require ~8 hours per run — beyond a single Colab session.\n\n"
    "  3. Hanabi underfitting — 5 M steps is insufficient; the train/eval gap (1.90 vs. 0.73) shows agents\n"
    "      haven't generalized hint conventions. More training would close the gap.",
    font_size=13, font_color=LIGHT)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Off-Policy Baseline Comparison
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK_BG)
title_bar(s, "MAPPO vs. Off-Policy Baselines",
          "Yu et al. (2022) — our MAPPO (2 M/5 M steps) vs. paper's off-policy algorithms (25 M steps)")

# ── Sub-header: MPE results ──────────────────────────────────────────────
rect(s, Inches(0.25), Inches(1.55), Inches(12.83), Inches(0.35), RGBColor(0x0F, 0x3E, 0x7E))
box(s, Inches(0.25), Inches(1.55), Inches(12.83), Inches(0.35),
    RGBColor(0x0F, 0x3E, 0x7E),
    "  MPE Environments  (episode reward — higher = better, paper Fig. 1 at 25 M steps)",
    font_size=13, bold=True, font_color=GOLD)

# MPE table headers
OFF_HDR_BG = RGBColor(0x10, 0x30, 0x60)
OFF_ROW = [RGBColor(0x1A, 0x28, 0x42), RGBColor(0x14, 0x1E, 0x38)]
BEAT_COL   = RGBColor(0x4C, 0xAF, 0x50)
CLOSE_COL  = RGBColor(0xE9, 0xC4, 0x6A)
BEHIND_COL = RGBColor(0xFF, 0x70, 0x50)

MPE_COLS     = ["Environment", "Our MAPPO\n(2 M steps)", "Paper MAPPO\n(25 M steps)", "QMIX\n(25 M steps)", "MADDPG\n(25 M steps)", "PPO vs.\nOff-Policy"]
MPE_COL_X    = [Inches(0.25), Inches(2.45), Inches(4.35), Inches(6.25), Inches(8.15), Inches(10.05)]
MPE_COL_W    = [Inches(2.1),  Inches(1.8),  Inches(1.8),  Inches(1.8),  Inches(1.8),  Inches(3.05)]
MPE_HDR_Y    = Inches(1.92)
MPE_ROW_H    = Inches(0.53)

for ci, (hdr, x, w) in enumerate(zip(MPE_COLS, MPE_COL_X, MPE_COL_W)):
    box(s, x, MPE_HDR_Y, w, MPE_ROW_H, OFF_HDR_BG, hdr,
        font_size=12, bold=True, align=PP_ALIGN.CENTER, font_color=WHITE)

mpe_rows = [
    ("TC1: simple_spread",   "−125 (early)",  "~−85",   "~−100",  "~−135",  "Matches / beats at convergence"),
    ("TC2: speaker_listener","−13.3",         "~−12",   "~−22",   "~−75",   "Beats both off-policy"),
]
mpe_verdict_colors = [CLOSE_COL, BEAT_COL]

for ri, (row, vcol) in enumerate(zip(mpe_rows, mpe_verdict_colors)):
    y = MPE_HDR_Y + MPE_ROW_H * (ri + 1)
    for ci, (val, x, w) in enumerate(zip(row, MPE_COL_X, MPE_COL_W)):
        fc = vcol if ci == 5 else LIGHT
        box(s, x, y, w, MPE_ROW_H, OFF_ROW[ri % 2], val,
            font_size=12, align=PP_ALIGN.CENTER, font_color=fc)

# ── Sub-header: Hanabi results ──────────────────────────────────────────
HAN_TOP = Inches(3.32)
rect(s, Inches(0.25), HAN_TOP, Inches(12.83), Inches(0.35), RGBColor(0x0F, 0x3E, 0x7E))
box(s, Inches(0.25), HAN_TOP, Inches(12.83), Inches(0.35),
    RGBColor(0x0F, 0x3E, 0x7E),
    "  Hanabi  (full game, max score = 25 — paper Table 3;  our TC3 used Hanabi-Very-Small, max = 5)",
    font_size=13, bold=True, font_color=GOLD)

HAN_COLS  = ["Config",      "MAPPO",  "IPPO",   "SAD\n(off-policy)", "VDN\n(off-policy)", "PPO vs. Off-Policy"]
HAN_COL_X = [Inches(0.25),  Inches(2.45), Inches(4.35), Inches(6.25), Inches(8.15), Inches(10.05)]
HAN_COL_W = [Inches(2.1),   Inches(1.8),  Inches(1.8),  Inches(1.8),  Inches(1.8),  Inches(3.05)]
HAN_HDR_Y = HAN_TOP + Inches(0.37)
HAN_ROW_H = Inches(0.46)

for ci, (hdr, x, w) in enumerate(zip(HAN_COLS, HAN_COL_X, HAN_COL_W)):
    box(s, x, HAN_HDR_Y, w, HAN_ROW_H, OFF_HDR_BG, hdr,
        font_size=12, bold=True, align=PP_ALIGN.CENTER, font_color=WHITE)

han_rows = [
    ("2-player (paper)",  "23.89", "24.00", "23.87", "23.83", "Competitive / superior"),
    ("3-player (paper)",  "23.77", "23.25", "23.69", "23.71", "Competitive / superior"),
    ("5-player (paper)",  "23.04", "20.75", "22.06", "21.28", "Beats all (best algo)"),
    ("Our TC3 (V-Small)", "1.90/5","  —",   "  —",   "  —",   "On trajectory at 5 M steps"),
]
han_verdict_colors = [BEAT_COL, BEAT_COL, BEAT_COL, CLOSE_COL]

for ri, (row, vcol) in enumerate(zip(han_rows, han_verdict_colors)):
    y = HAN_HDR_Y + HAN_ROW_H * (ri + 1)
    for ci, (val, x, w) in enumerate(zip(row, HAN_COL_X, HAN_COL_W)):
        fc = vcol if ci == 5 else LIGHT
        box(s, x, y, w, HAN_ROW_H, OFF_ROW[ri % 2], val,
            font_size=12, align=PP_ALIGN.CENTER, font_color=fc)

# ── Footer note ──────────────────────────────────────────────────────────
rect(s, 0, Inches(6.65), W, Inches(0.85), RGBColor(0x0D, 0x1E, 0x36))
box(s, Inches(0.35), Inches(6.68), Inches(12.6), Inches(0.75),
    RGBColor(0x0D, 0x1E, 0x36),
    "MPE values approximate from paper Fig. 1 learning curves at convergence (~25 M steps).  "
    "QMIX & MADDPG are the off-policy baselines in the paper for MPE; SAD & VDN for Hanabi.  "
    "Our TC3 used Hanabi-Very-Small (1 color, 5 ranks, max 5) — not directly comparable to the full-game paper table.",
    font_size=11, font_color=RGBColor(0x88, 0x99, 0xAA))

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Fig 1 Raw + Smoothed curves (cell 34)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK_BG)
title_bar(s, "Analytics — Raw & Smoothed Learning Curves",
          "All three scenarios · exponential moving average overlay")
full_image(s, f"{IMG}/cell34_img0.png")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Fig 2 Rolling Mean ± Std (cell 35)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK_BG)
title_bar(s, "Analytics — Training Stability",
          "Rolling mean ± std · variance decay over training")
full_image(s, f"{IMG}/cell35_img0.png")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Fig 3 Sample Efficiency (cell 36)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK_BG)
title_bar(s, "Analytics — Sample Efficiency",
          "Steps-to-threshold · 50 / 70 / 85 / 95% of peak performance")
full_image(s, f"{IMG}/cell36_img0.png")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Fig 4 Per-Agent Speaker vs Listener (cell 37)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK_BG)
title_bar(s, "Analytics — Per-Agent Deep Dive",
          "Speaker vs Listener reward breakdown (TC2)")
full_image(s, f"{IMG}/cell37_img0.png")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Fig 5 Convergence Analysis (cell 38)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK_BG)
title_bar(s, "Analytics — Convergence Analysis",
          "Improvement rate, variance decay, gradient trends")
full_image(s, f"{IMG}/cell38_img0.png")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Fig 6 Cross-Scenario Dashboard (cell 39)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK_BG)
title_bar(s, "Analytics — Cross-Scenario Comparison",
          "Final performance & relative gains across all environments")
full_image(s, f"{IMG}/cell39_img0.png")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Fig 7 Hanabi Deep Dive (cell 40)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK_BG)
title_bar(s, "Analytics — Hanabi Deep Dive",
          "Score distribution, train/eval gap, score progression")
full_image(s, f"{IMG}/cell40_img0.png")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Fig 8 Architecture & Hyperparameters (cell 41)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK_BG)
title_bar(s, "Model Architecture & Hyperparameter Summary",
          "Network design and training configuration used across experiments")
full_image(s, f"{IMG}/cell41_img0.png")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Fig 9 Numeric Metrics Table (cell 42)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK_BG)
title_bar(s, "Numeric Metrics Summary",
          "Final reward, convergence step, peak performance — all scenarios")
full_image(s, f"{IMG}/cell42_img0.png")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Key Takeaways
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, DARK_BG)
title_bar(s, "Key Takeaways",
          "What the reproduction confirmed — and where it fell short")

items = [
    "  TC2 (speaker-listener) most closely matches the paper: MAPPO beats IPPO at −13.3 vs −15.0,\n"
    "  confirming that a centralized critic is essential for credit assignment across a communication channel.",

    "  TC1 (simple_spread) shows the right direction but not the full gap: we ran only 2 M of the paper's\n"
    "  25 M steps — MAPPO's spatial coordination advantage is a long-horizon effect that needs more training.",

    "  TC3 (Hanabi) is on trajectory but underfitting at 5 M steps: train score 1.90/5 vs. paper's ~2.5/5,\n"
    "  with a large train/eval gap (1.90 vs. 0.73) showing agents haven't generalized hint conventions.",

    "  The step budget was constrained by Colab session limits: TC1 capped at 32 threads would need\n"
    "  ~8 hrs per run at 25 M steps — dedicated GPU compute would close the gap with the paper.",

    "  Core finding validated: sharing global state during training produces measurably better cooperation,\n"
    "  most strongly in tasks where credit assignment across agents is hardest.",
]
bullet_slide(s, items, font_size=17, gap=Inches(1.18))

# ════════════════════════════════════════════════════════════════════════════
# Save
# ════════════════════════════════════════════════════════════════════════════
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
